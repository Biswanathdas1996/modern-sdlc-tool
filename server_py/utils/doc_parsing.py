"""Multimodal document parsing utility.

Extracts text AND images from PDF, Word (DOCX), and PowerPoint (PPTX) documents.
Returns structured content with text blocks and extracted images for vision LLM captioning.

Supports:
  - PDF: PyMuPDF (fitz) for text + embedded image extraction
  - DOCX: python-docx for paragraphs + embedded images
  - PPTX: python-pptx for slide text + embedded images
  - TXT/MD/CSV/JSON: plain text (no image extraction)
"""

import io
import json
from typing import List, Dict, Any, Optional, Tuple
from dataclasses import dataclass, field
from core.logging import log_info, log_error


@dataclass
class ExtractedImage:
    image_bytes: bytes
    source_page: int
    source_type: str
    mime_type: str
    width: Optional[int] = None
    height: Optional[int] = None
    caption: Optional[str] = None


@dataclass
class ParsedContent:
    text_blocks: List[Dict[str, Any]] = field(default_factory=list)
    images: List[ExtractedImage] = field(default_factory=list)
    metadata: Dict[str, Any] = field(default_factory=dict)

    @property
    def full_text(self) -> str:
        return "\n\n".join(block["text"] for block in self.text_blocks if block.get("text"))

    @property
    def combined_text_with_captions(self) -> str:
        sections = []
        img_by_page: Dict[int, List[ExtractedImage]] = {}
        for img in self.images:
            if img.caption:
                img_by_page.setdefault(img.source_page, []).append(img)

        for block in self.text_blocks:
            sections.append(block["text"])
            page = block.get("page", -1)
            if page in img_by_page:
                for img in img_by_page[page]:
                    sections.append(f"[Image on {img.source_type} {img.source_page + 1}]: {img.caption}")
                del img_by_page[page]

        for page, imgs in sorted(img_by_page.items()):
            for img in imgs:
                sections.append(f"[Image on {img.source_type} {img.source_page + 1}]: {img.caption}")

        return "\n\n".join(s for s in sections if s.strip())


MIN_IMAGE_SIZE = 5000
MIN_IMAGE_DIMENSION = 50
MAX_IMAGES_PER_DOC = 30


def _is_meaningful_image(img_bytes: bytes, width: Optional[int] = None, height: Optional[int] = None) -> bool:
    if len(img_bytes) < MIN_IMAGE_SIZE:
        return False
    if width and height:
        if width < MIN_IMAGE_DIMENSION or height < MIN_IMAGE_DIMENSION:
            return False
    return True


def _normalize_image_to_png(img_bytes: bytes, mime_type: str) -> Tuple[bytes, str]:
    if mime_type in ("image/png", "image/jpeg", "image/jpg"):
        return img_bytes, mime_type
    try:
        import fitz
        pix = fitz.Pixmap(img_bytes)
        if pix.alpha:
            pix = fitz.Pixmap(fitz.csRGB, pix)
        png_bytes = pix.tobytes("png")
        return png_bytes, "image/png"
    except Exception:
        return img_bytes, mime_type


def parse_pdf(file_content: bytes) -> ParsedContent:
    import fitz

    result = ParsedContent()
    result.metadata["format"] = "pdf"

    try:
        doc = fitz.open(stream=file_content, filetype="pdf")
        result.metadata["page_count"] = len(doc)

        for page_idx in range(len(doc)):
            page = doc[page_idx]

            text = page.get_text("text")
            if text and text.strip():
                result.text_blocks.append({
                    "text": text.strip(),
                    "page": page_idx,
                    "source": f"page_{page_idx + 1}",
                })

            if len(result.images) >= MAX_IMAGES_PER_DOC:
                continue

            image_list = page.get_images(full=True)
            for img_info in image_list:
                try:
                    xref = img_info[0]
                    base_image = doc.extract_image(xref)
                    if not base_image:
                        continue

                    img_bytes = base_image["image"]
                    mime = f"image/{base_image.get('ext', 'png')}"
                    w = base_image.get("width", 0)
                    h = base_image.get("height", 0)

                    if not _is_meaningful_image(img_bytes, w, h):
                        continue

                    img_bytes, mime = _normalize_image_to_png(img_bytes, mime)

                    result.images.append(ExtractedImage(
                        image_bytes=img_bytes,
                        source_page=page_idx,
                        source_type="page",
                        mime_type=mime,
                        width=w,
                        height=h,
                    ))
                except Exception as e:
                    log_error(f"Error extracting image from PDF page {page_idx}", "doc_parsing", e)

        doc.close()
        log_info(
            f"PDF parsed: {len(result.text_blocks)} text blocks, {len(result.images)} images from {result.metadata['page_count']} pages",
            "doc_parsing"
        )

    except Exception as e:
        log_error("Error parsing PDF", "doc_parsing", e)
        try:
            from PyPDF2 import PdfReader
            reader = PdfReader(io.BytesIO(file_content))
            for i, page in enumerate(reader.pages):
                text = page.extract_text() or ""
                if text.strip():
                    result.text_blocks.append({
                        "text": text.strip(),
                        "page": i,
                        "source": f"page_{i + 1}",
                    })
            result.metadata["page_count"] = len(reader.pages)
            result.metadata["fallback"] = "PyPDF2"
            log_info(f"PDF parsed with PyPDF2 fallback: {len(result.text_blocks)} text blocks", "doc_parsing")
        except Exception as e2:
            log_error("PyPDF2 fallback also failed", "doc_parsing", e2)

    return result


def parse_docx(file_content: bytes) -> ParsedContent:
    from docx import Document

    result = ParsedContent()
    result.metadata["format"] = "docx"

    try:
        doc = Document(io.BytesIO(file_content))

        for i, para in enumerate(doc.paragraphs):
            if para.text.strip():
                result.text_blocks.append({
                    "text": para.text.strip(),
                    "page": i // 40,
                    "source": f"paragraph_{i + 1}",
                })

        try:
            for rel in doc.part.rels.values():
                if "image" in rel.reltype:
                    try:
                        img_blob = rel.target_part.blob
                        content_type = rel.target_part.content_type or "image/png"

                        if not _is_meaningful_image(img_blob):
                            continue

                        if len(result.images) >= MAX_IMAGES_PER_DOC:
                            break

                        result.images.append(ExtractedImage(
                            image_bytes=img_blob,
                            source_page=0,
                            source_type="document",
                            mime_type=content_type,
                        ))
                    except Exception as e:
                        log_error("Error extracting DOCX image", "doc_parsing", e)
        except Exception as e:
            log_error("Error accessing DOCX relationships", "doc_parsing", e)

        log_info(
            f"DOCX parsed: {len(result.text_blocks)} text blocks, {len(result.images)} images",
            "doc_parsing"
        )

    except Exception as e:
        log_error("Error parsing DOCX", "doc_parsing", e)

    return result


def parse_pptx(file_content: bytes) -> ParsedContent:
    from pptx import Presentation
    from pptx.enum.shapes import MSO_SHAPE_TYPE

    result = ParsedContent()
    result.metadata["format"] = "pptx"

    try:
        prs = Presentation(io.BytesIO(file_content))
        result.metadata["slide_count"] = len(prs.slides)

        for slide_idx, slide in enumerate(prs.slides):
            slide_texts = []

            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        text = para.text.strip()
                        if text:
                            slide_texts.append(text)

                if shape.has_table:
                    table = shape.table
                    for row in table.rows:
                        row_texts = [cell.text.strip() for cell in row.cells if cell.text.strip()]
                        if row_texts:
                            slide_texts.append(" | ".join(row_texts))

                if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                    if len(result.images) >= MAX_IMAGES_PER_DOC:
                        continue
                    try:
                        image = shape.image
                        img_blob = image.blob
                        content_type = image.content_type or "image/png"

                        if not _is_meaningful_image(img_blob):
                            continue

                        result.images.append(ExtractedImage(
                            image_bytes=img_blob,
                            source_page=slide_idx,
                            source_type="slide",
                            mime_type=content_type,
                            width=shape.width,
                            height=shape.height,
                        ))
                    except Exception as e:
                        log_error(f"Error extracting PPTX image from slide {slide_idx}", "doc_parsing", e)

                if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                    try:
                        for child in shape.shapes:
                            if child.shape_type == MSO_SHAPE_TYPE.PICTURE:
                                if len(result.images) >= MAX_IMAGES_PER_DOC:
                                    break
                                image = child.image
                                img_blob = image.blob
                                content_type = image.content_type or "image/png"
                                if _is_meaningful_image(img_blob):
                                    result.images.append(ExtractedImage(
                                        image_bytes=img_blob,
                                        source_page=slide_idx,
                                        source_type="slide",
                                        mime_type=content_type,
                                    ))
                    except Exception:
                        pass

            if slide_texts:
                result.text_blocks.append({
                    "text": "\n".join(slide_texts),
                    "page": slide_idx,
                    "source": f"slide_{slide_idx + 1}",
                })

        log_info(
            f"PPTX parsed: {len(result.text_blocks)} text blocks, {len(result.images)} images from {result.metadata.get('slide_count', 0)} slides",
            "doc_parsing"
        )

    except Exception as e:
        log_error("Error parsing PPTX", "doc_parsing", e)

    return result


def parse_plain_text(file_content: bytes, content_type: str = "text/plain") -> ParsedContent:
    result = ParsedContent()

    try:
        if content_type == "application/json":
            json_content = json.loads(file_content.decode("utf-8"))
            text = json.dumps(json_content, indent=2)
        else:
            text = file_content.decode("utf-8", errors="ignore")

        if text.strip():
            result.text_blocks.append({
                "text": text.strip(),
                "page": 0,
                "source": "full_document",
            })

        result.metadata["format"] = content_type.split("/")[-1]

    except Exception as e:
        log_error("Error parsing text file", "doc_parsing", e)
        text = file_content.decode("utf-8", errors="ignore")
        if text.strip():
            result.text_blocks.append({
                "text": text.strip(),
                "page": 0,
                "source": "full_document",
            })

    return result


CONTENT_TYPE_MAP = {
    "application/pdf": parse_pdf,
    "application/vnd.openxmlformats-officedocument.wordprocessingml.document": parse_docx,
    "application/vnd.openxmlformats-officedocument.presentationml.presentation": parse_pptx,
}

EXTENSION_MAP = {
    ".pdf": parse_pdf,
    ".docx": parse_docx,
    ".pptx": parse_pptx,
}

TEXT_CONTENT_TYPES = {
    "text/plain", "text/markdown", "text/csv", "application/json",
}


def parse_document(file_content: bytes, content_type: str, filename: str = "") -> ParsedContent:
    if content_type in CONTENT_TYPE_MAP:
        return CONTENT_TYPE_MAP[content_type](file_content)

    ext = ""
    if filename:
        ext = "." + filename.rsplit(".", 1)[-1].lower() if "." in filename else ""

    if ext in EXTENSION_MAP:
        return EXTENSION_MAP[ext](file_content)

    if content_type in TEXT_CONTENT_TYPES or ext in {".txt", ".md", ".csv", ".json"}:
        return parse_plain_text(file_content, content_type)

    return parse_plain_text(file_content, content_type)
